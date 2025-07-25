#!/usr/bin/env python3
"""
Enhanced ChatGPT Canvas to PowerPoint Converter
Addresses all potential failure points with robust error handling and enhanced features.
"""

import os
import re
import sys
import time
import json
import logging
import tempfile
import unicodedata
from pathlib import Path
from urllib.parse import urlparse
from typing import Optional, Tuple, List, Dict, Any
from dataclasses import dataclass
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from docx import Document
from pptx.util import Cm
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from bs4.element import Tag
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from bs4 import BeautifulSoup, NavigableString

# Core dependencies with error handling

# === CONFIGURATION ===
@dataclass
class Config:
    """Configuration settings for the converter"""
    max_wait_time: int = 30
    page_load_timeout: int = 45
    retry_attempts: int = 3
    max_slide_content_length: int = 1000
    max_filename_length: int = 100
    supported_image_formats: List[str] = None
    font_fallbacks: Dict[str, str] = None

    def __post_init__(self):
        if self.supported_image_formats is None:
            self.supported_image_formats = ['.png', '.jpg', '.jpeg', '.gif', '.bmp']
        if self.font_fallbacks is None:
            self.font_fallbacks = {
                'default': 'Calibri',
                'code': 'Courier New',
                'math': 'Cambria Math',
                'fallback': 'Arial'
            }


# === LOGGING SETUP ===
def setup_logging(log_level: str = 'INFO') -> logging.Logger:
    """Setup comprehensive logging"""
    logger = logging.getLogger('canvas_converter')
    logger.setLevel(getattr(logging, log_level.upper()))

    if not logger.handlers:
        handler = logging.StreamHandler()
        formatter = logging.Formatter(
            '%(asctime)s - %(name)s - %(levelname)s - %(message)s'
        )
        handler.setFormatter(formatter)
        logger.addHandler(handler)

    return logger


# === UTILITIES ===
class SafeFilename:
    """Utility class for safe filename handling"""

    @staticmethod
    def sanitize(filename: str, max_length: int = 100) -> str:
        """Create safe filename with proper sanitization"""
        # Remove or replace dangerous characters
        filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
        filename = re.sub(r'\s+', '_', filename)

        # Handle Unicode characters
        filename = unicodedata.normalize('NFKD', filename)
        filename = ''.join(c for c in filename if ord(c) < 128)

        # Truncate and ensure extension
        if len(filename) > max_length - 5:  # Leave room for .pptx
            filename = filename[:max_length - 5]

        return filename if filename else 'canvas_export'

    @staticmethod
    def ensure_unique(filepath: Path) -> Path:
        """Ensure filename is unique by adding counter if needed"""
        if not filepath.exists():
            return filepath

        counter = 1
        stem = filepath.stem
        suffix = filepath.suffix
        parent = filepath.parent

        while True:
            new_path = parent / f"{stem}_{counter}{suffix}"
            if not new_path.exists():
                return new_path
            counter += 1




# class TableProcessor:
#     """Enhanced table processing with error handling"""

#     def __init__(self, logger: logging.Logger):
#         self.logger = logger

#     def clean_and_validate_tables(self, soup: BeautifulSoup, content_div: Tag) -> Tag:
#         """Clean tables with comprehensive error handling"""
#         try:
#             tables = content_div.find_all("table")
#             self.logger.info(f"Processing {len(tables)} tables")

#             for i, table in enumerate(tables):
#                 try:
#                     self._clean_single_table(soup, table, i)
#                 except Exception as e:
#                     self.logger.warning(f"Failed to clean table {i}: {e}")
#                     # Remove problematic table rather than crash
#                     table.decompose()

#             return content_div
#         except Exception as e:
#             self.logger.error(f"Table processing failed: {e}")
#             return content_div

#     def _clean_single_table(self, soup: BeautifulSoup, table: Tag, table_index: int) -> None:
#         """Clean a single table with validation"""
#         rows = table.find_all("tr")
#         if not rows:
#             self.logger.warning(f"Table {table_index} has no rows")
#             return

#         # Calculate maximum columns safely
#         max_cols = 0
#         valid_rows = []

#         for row in rows:
#             cells = row.find_all(["td", "th"])
#             if cells:  # Only process rows with cells
#                 max_cols = max(max_cols, len(cells))
#                 valid_rows.append((row, cells))

#         if max_cols == 0:
#             self.logger.warning(f"Table {table_index} has no valid cells")
#             return

#         # Normalize all rows to have the same number of columns
#         for row, cells in valid_rows:
#             while len(cells) < max_cols:
#                 new_cell = soup.new_tag("td")
#                 new_cell.string = ""
#                 row.append(new_cell)
#                 cells.append(new_cell)


class PowerPointGenerator:
    """Enhanced PowerPoint generation with advanced features"""

    def __init__(self, config: Config, logger: logging.Logger):
        self.speaker_notes_txt = [] 
        self.notes_seen = set()
        self.config = config
        self.logger = logger
        self.slide_count = 0
        self.watermark_path = Path("assets/watermark.jpg")  # use relative path


    def create_enhanced_presentation(self, content_div: Tag, output_path: Path, title: str = None) -> bool:
        """Create PowerPoint with enhanced features and error handling"""
        try:
            prs = Presentation()

            self._set_default_fonts(prs)

            # Add title slide
            # Manually extract first heading/subheading/speaker notes
            elements = content_div.find_all(["h1", "h2", "p"], recursive=True)
            for el in elements:
                if el.name == "p" and re.search(r"speaker notes\s*:", el.get_text(), flags=re.IGNORECASE):
                    break  # Prevents same speaker note being parsed again later
            heading = None
            subheading = None
            speaker_notes = None

            for i, el in enumerate(elements):
                text = el.get_text(strip=True)

                # Remove "Slide x:" prefix if present
                text = re.sub(r'^slide\s*\d+\s*:\s*', '', text, flags=re.IGNORECASE)

                # Detect speaker notes first
                if re.search(r"speaker notes\s*:\s*", text, flags=re.IGNORECASE):
                    _, notes = re.split(r"speaker notes\s*:\s*", text, flags=re.IGNORECASE, maxsplit=1)
                    speaker_notes = notes.strip()
                    el.decompose()
                    continue

                # Assign heading and subheading next
                if not heading:
                    heading = text
                    elements[i].decompose()
                elif not subheading:
                    subheading = text
                    elements[i].decompose()

                # Once both are found, break
                if heading and subheading and speaker_notes:
                    break

            # Add custom title slide
            self.add_custom_title_slide(prs, heading or " ", subheading or " ", speaker_notes or "")


            # Process content elements
            self._process_content_elements(prs, content_div)

            # Ensure we have at least one slide
            if len(prs.slides) == 0:
                self._add_fallback_slide(prs, "No Content Found", "The canvas appears to be empty or could not be processed.")

            # Save with error handling
            self._save_presentation(prs, output_path)

            if self.speaker_notes_txt:
                self._save_speaker_notes_textfile(output_path, self.speaker_notes_txt)

            self.logger.info(f"✅ PowerPoint created with {len(prs.slides)} slides: {output_path}")

            return True

        except Exception as e:
            self.logger.error(f"PowerPoint generation failed: {e}")
            return False

    def _add_watermark(self, slide):
        """Insert watermark image on a slide"""
        try:
            if self.watermark_path.is_file():
                slide.shapes.add_picture(
                    str(self.watermark_path),
                    Inches(8.016),   # Horizontal position (unchanged)
                    Inches(-0.012),  # Vertical position (unchanged)
                    width=Cm(4.31),
                    height=Cm(1.87)
                )
            else:
                self.logger.warning(f"Watermark image not found: {self.watermark_path}")
        except Exception as e:
            self.logger.warning(f"Failed to insert watermark: {e}")

    def add_custom_title_slide(self, prs: Presentation, heading: str, subheading: str, speaker_notes: str) -> None:
        try:
            title_slide_layout = prs.slide_layouts[0]  # Title Slide layout
            slide = prs.slides.add_slide(title_slide_layout)

            slide.shapes.title.text = heading or ""
            clean_subheading = re.sub(r'^subtitle\s*:\s*', '', subheading, flags=re.IGNORECASE).strip()
            if len(slide.placeholders) > 1:
                 slide.placeholders[1].text = clean_subheading or ""

            # Add speaker notes
            if speaker_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = speaker_notes.strip()
                slide_index = prs.slides.index(slide)
                notes_key = (slide_index + 1, speaker_notes.strip())  # 1-based slide number
                if notes_key not in self.notes_seen and speaker_notes.strip():
                   self.notes_seen.add(notes_key)
                   self.speaker_notes_txt.append(notes_key)
            self._add_watermark(slide) 
            self.slide_count += 1
        except Exception as e:
            self.logger.warning(f"Failed to add custom title slide: {e}")



    def _process_content_elements(self, prs: Presentation, content_div: Tag) -> None:
        """Process all content elements with enhanced handling"""
        current_slide = None
        content_box = None

        # Get all relevant elements in document order
        elements = content_div.find_all([
            "h1", "h2", "h3", "h4", "h5", "h6",
            "p", "ul", "ol", "table", "pre", "code",
            "blockquote", "img", "span","div"
        ], recursive=True)

        code_buffer = []
        processed_elements = set()  # Track processed elements to avoid duplicates

        for element in elements:
            # Skip if already processed
            element_id = id(element)
            if element_id in processed_elements:
                continue

            element_text = element.get_text(strip=True)
            print(f"[DEBUG] Element text: {repr(element_text)}")

            # Handle speaker notes
            # Handle speaker notes
            match = re.search(r"speaker notes\s*:\s*", element_text, flags=re.IGNORECASE)
            if match:
                content_part, notes_part = re.split(r"speaker notes\s*:\s*", element_text, flags=re.IGNORECASE, maxsplit=1)

                # If we have a current slide, add notes to it
                if current_slide is not None:
                    slide = current_slide
                    notes_slide = slide.notes_slide
                    notes_key = (self.slide_count, notes_part.strip())
                    if notes_part.strip() and notes_key not in self.notes_seen:
                           self.notes_seen.add(notes_key)
                           self.speaker_notes_txt.append(notes_key)
                           notes_slide.notes_text_frame.text = notes_part.strip()

                # Mark as processed and skip this element entirely
                processed_elements.add(element_id)
                continue  # Skip the rest of the processing for this element

            # Handle consecutive cm-line blocks as one code block
            if element.name == "div" and "cm-line" in element.get("class", []):
                print(f"[DEBUG] Detected cm-line: {element.get_text(strip=True)}")
                code_text = element.get_text(strip=True)
                if code_text:
                    code_buffer.append(code_text)
                processed_elements.add(element_id)
                continue

            # If the current element is NOT a cm-line AND we have code buffered, flush it
            if code_buffer:
                current_slide, content_box = self._ensure_slide(prs, current_slide, "Content")
                self._add_code_content(content_box, "\n".join(code_buffer))
                code_buffer = []

            # Skip elements that are part of other elements to avoid duplication
            if element.name in ["p", "span"] and element.find_parents(["ul", "ol", "li"]):
                processed_elements.add(element_id)
                continue

            if element.name == "p" and element.find_parent("li"):
                processed_elements.add(element_id)
                continue

            if element.name == "span" and (element.find_parent("p") or element.find_parent("li")):
                processed_elements.add(element_id)
                continue

            # Skip <code> if it's inside a <pre>
            if element.name == "code" and element.find_parent("pre"):
                processed_elements.add(element_id)
                continue

            try:
                element_type = element.name

                # Handle headings - create new slides
                if element_type in ["h1", "h2", "h3", "h4", "h5", "h6"]:
                    current_slide, content_box = self._add_content_slide(prs, element.get_text(strip=True))

                # Handle lists
                elif element_type in ["ol", "ul"] and not element.find_parent(["ul", "ol"]):
                    current_slide, content_box = self._ensure_slide(prs, current_slide, "List")
                    self._add_list_content(content_box, element)
                    processed_elements.add(element_id)
                    continue

                # elif element.name == "p":
                #    current_slide, content_box = self._handle_paragraph_element(prs, element, current_slide, content_box)
                #    continue

                # Handle tables
                elif element_type == "table":
                   current_slide, content_box = self._ensure_slide(prs, current_slide, "Content")
                   self._add_table_to_slide(current_slide, element)
                   processed_elements.add(element_id)
                   continue

                # Handle code blocks
                elif element_type in ["pre", "code"]:
                    current_slide, content_box = self._ensure_slide(prs, current_slide, "Code")
                    self._add_code_content(content_box, element)
                    continue
                # Mark as processed
                processed_elements.add(element_id)

            except Exception as e:
                self.logger.warning(f"Failed to process element {element_type}: {e}")
                processed_elements.add(element_id)
                continue

            # Handle any remaining code buffer
            if code_buffer:
                current_slide, content_box = self._ensure_slide(prs, current_slide, "Content")
                self._add_code_content(content_box, "\n".join(code_buffer))

    def _add_content_slide(self, prs: Presentation, title: str) -> Tuple[Any, Any]:
        """Add a new content slide and apply custom formatting to the body placeholder"""
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        # Clean the title and set it
        clean_title = re.sub(r'^slide\s*\d+\s*:\s*', '', title, flags=re.IGNORECASE)
        slide.shapes.title.text = clean_title[:100] + "..." if len(clean_title) > 100 else clean_title
        # Set heading font size to 38pt
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.size = Pt(38)

        # Get the body placeholder
        content_box = slide.placeholders[1]
        text_frame = content_box.text_frame
        text_frame.clear()

        # Apply new custom shape settings as requested
        content_box.left = Cm(1.27)
        content_box.top = Cm(5.28)
        content_box.width = Cm(21.78)
        content_box.height = Cm(8.08)

        text_frame.word_wrap = True
        # Enable auto shrink to fit text
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        text_frame.margin_left = Cm(0.25)
        text_frame.margin_right = Cm(0.25)
        text_frame.margin_top = Cm(0.13)
        text_frame.margin_bottom = Cm(0.13)
        self._add_watermark(slide)
        self.slide_count += 1
        return slide, text_frame

    def _ensure_slide(self, prs: Presentation, current_slide: Any, default_title: str) -> Tuple[Any, Any]:
        """Ensure we have a slide to work with"""
        if current_slide is None:
            return self._add_content_slide(prs, default_title)

        # Return existing slide and its content box
        content_box = current_slide.placeholders[1]
        return current_slide, content_box

    def _handle_paragraph_element(self, prs: Presentation, element: Tag, current_slide: Any, content_box: Any) -> Tuple[Any, Any]:
        """Handle standalone <p> elements not part of lists"""
        if not element.get_text(strip=True):
            return current_slide, content_box

        current_slide, content_box = self._ensure_slide(prs, current_slide, "Content")
        self._add_paragraph_content(content_box, element)
        return current_slide, content_box

    def _add_paragraph_content(self, content_box: Any, element: Tag) -> None:
        """Add paragraph with smart formatting"""
        print("Para content being called")
        text = element.get_text(strip=True)
        if not text or len(text) > self.config.max_slide_content_length:
            return
        if "speaker notes:" in text.lower():
           content_part, notes_part = re.split(r'speaker notes\s*:\s*', text, flags=re.IGNORECASE, maxsplit=1)
           text = content_part.strip()

           # Add speaker notes to the slide
           notes_slide = content_box.part.slide.notes_slide
           notes_slide.notes_text_frame.text = notes_part.strip()

        # Check if this looks like a bullet point
        bullet_patterns = [r'^\s*[-•·]\s+', r'^\s*\d+\.\s+', r'^\s*[a-zA-Z]\.\s+']
        is_bullet = any(re.match(pattern, text) for pattern in bullet_patterns)

        para = content_box.text_frame.add_paragraph()
        para.text = text
        para.level = 1 if is_bullet else 0


        self._set_font_safely(para, text, 'default')

    def _add_list_content(self, content_box: Any, list_element: Tag) -> None:
        """Add list with proper hierarchy"""
        self._process_list_recursive(content_box, list_element, 0)

    def _process_list_recursive(self, content_box: Any, list_element: Tag, level: int) -> None:
        """Process lists recursively with proper nesting"""
        max_level = 4  # PowerPoint limitation
        print("rec list being called")
        if len(content_box.text_frame.paragraphs) == 1 and not content_box.text_frame.paragraphs[0].text.strip():
            p = content_box.text_frame.paragraphs[0]
            content_box.text_frame._element.remove(p._element)
        for li in list_element.find_all("li", recursive=False):
            try:
                # Get text content, excluding nested lists
                text_parts = []
                for item in li.children:
                    if isinstance(item, NavigableString):
                        text_parts.append(str(item).strip())
                    elif isinstance(item, Tag) and item.name not in ["ul", "ol"]:
                        text_parts.append(item.get_text(" ", strip=True))

                text = " ".join(text_parts).strip()
                if text:
                    if len(content_box.text_frame.paragraphs) == 1 and not content_box.text_frame.paragraphs[0].text.strip():
                        para = content_box.text_frame.paragraphs[0]
                    else:
                        para = content_box.text_frame.add_paragraph()
                    para.text = text
                    para.level = min(level, max_level)
                    para.space_before = Pt(0)
                    para.space_after = Pt(0)
                    para.margin_left = Pt(0)
                    para.left_indent = Pt(0)

                    self._set_font_safely(para, text, 'default')

                # Process nested lists
                nested_lists = li.find_all(["ul", "ol"], recursive=False)
                for nested in nested_lists:
                    self._process_list_recursive(content_box, nested, level + 1)

            except Exception as e:
                self.logger.debug(f"List item processing failed: {e}")
        # ✅ Trigger shrink-to-fit for bullet list after all items are added
        content_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    def _add_table_to_slide(self, slide: Any, table_element: Tag) -> None:
        """Insert table into an existing slide"""
        try:
            rows = table_element.find_all("tr")
            if not rows:
                return

            # Calculate table dimensions
            max_cols = max(len(row.find_all(["td", "th"])) for row in rows if row.find_all(["td", "th"]))
            num_rows = len(rows)

            if max_cols == 0 or num_rows == 0:
                return

            # Table dimensions and positioning
            left = Inches(0.5)
            top = Inches(1.77)  # Position below title or existing content
            width = Inches(9)
            height = Inches(min(5.5, 0.5 + 0.4 * num_rows))

            # Create table
            table_shape = slide.shapes.add_table(num_rows, max_cols, left, top, width, height)
            table = table_shape.table

            for i, row in enumerate(rows):
                cells = row.find_all(["td", "th"])
                is_header = any(cell.name == "th" for cell in cells)

                for j in range(max_cols):
                    cell_text = ""
                    if j < len(cells):
                        cell_text = cells[j].get_text(strip=True)

                    cell = table.cell(i, j)
                    cell.text = cell_text[:200]

                    if is_header and i == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                run.font.bold = True
        except Exception as e:
            self.logger.warning(f"Table insertion failed: {e}")


    def _add_code_content(self, content_box: Any, code_text: str) -> None:
        """Add code with monospace formatting and no bullets"""
        if not code_text.strip():
            return

        para = content_box.text_frame.add_paragraph()
        para.text = code_text[:self.config.max_slide_content_length]
        para.level = 0  # Make sure it's top-level

        # 🔧 Safely remove bullets
        pPr = para._element.get_or_add_pPr()
        for bullet_tag in ['a:buAutoNum', 'a:buChar', 'a:buNone']:
            tag = pPr.find(qn(bullet_tag))
            if tag is not None:
                pPr.remove(tag)

        # ⛑️ Add bullet=None explicitly (if using a style that enforces bullets)
        buNone = OxmlElement('a:buNone')
        pPr.append(buNone)

        self._set_font_safely(para, code_text, 'code')


    def _add_quote_content(self, content_box: Any, element: Tag) -> None:
        """Add blockquote with special formatting"""
        quote_text = element.get_text(strip=True)
        if not quote_text:
            return

        para = content_box.text_frame.add_paragraph()
        para.text = f'"{quote_text}"'
        para.level = 0

        self._set_font_safely(para, quote_text, 'default')
        try:
            para.font.italic = True
        except Exception:
            pass

    # '''def _add_math_content(self, content_box: Any, element: Tag) -> None:
    #     """Add mathematical expressions"""
    #     math_text = element.get_text(strip=True)
    #     if not math_text:
    #         return

    #     para = content_box.text_frame.add_paragraph()
    #     para.text = f"Formula: {math_text}"

    #     try:
    #         para.font.name = self.config.font_fallbacks['math']
    #         para.font.size = Pt(26)
    #     except Exception:
    #         pass

    # def _is_math_element(self, element: Tag) -> bool:
    #     """Check if element contains mathematical expressions"""
    #     class_names = element.get("class", [])
    #     return any("katex" in str(cls).lower() or "math" in str(cls).lower() for cls in class_names)'''

    # def _get_appropriate_font(self, text: str) -> str:
    #     """Get appropriate font based on text content"""
    #     # Check for non-ASCII characters (might need special font handling)
    #     if any(ord(c) > 127 for c in text):
    #         return self.config.font_fallbacks['fallback']

    #     # Check for code-like content
    #     if re.search(r'[{}();=<>]', text) and len(text.split()) < 10:
    #         return self.config.font_fallbacks['code']

    #     return self.config.font_fallbacks['default']

    def _add_fallback_slide(self, prs: Presentation, title: str, content: str) -> None:
        """Add fallback slide when no content is found"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
        self.slide_count += 1

    def _set_font_safely(self, paragraph, text_content, font_type='default'):
        """Safely set font with proper error handling and logging"""
        try:
            if font_type == 'code':
                font_name = self.config.font_fallbacks['code']
                font_size = Pt(20)
            elif font_type == 'heading':
                font_name = self.config.font_fallbacks['calibri']
                font_size = Pt(28)
            else:
                font_name = self._get_appropriate_font(text_content)
                font_size = Pt(24)

            paragraph.font.name = font_name
            paragraph.font.size = font_size

            self.logger.debug(f"Font set successfully: {font_name}, {font_size}")

        except Exception as e:
            self.logger.warning(f"Font setting failed for '{text_content[:50]}...': {e}")
            # Try fallback
            try:
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(22)
            except Exception as fallback_error:
                self.logger.error(f"Even fallback font failed: {fallback_error}")

    def _set_default_fonts(self, prs: Presentation):
        """Set consistent default fonts across all slide layouts"""
        try:
            for layout in prs.slide_layouts:
                for placeholder in layout.placeholders:
                    if hasattr(placeholder, 'text_frame'):
                        for paragraph in placeholder.text_frame.paragraphs:
                            try:
                                paragraph.font.name = self.config.font_fallbacks['default']
                                paragraph.font.size = Pt(22)
                            except:
                                continue
        except Exception as e:
            self.logger.debug(f"Default font setting failed: {e}")

    def _save_presentation(self, prs: Presentation, output_path: Path) -> None:
        """Save presentation with comprehensive error handling"""
        try:
            # Ensure directory exists
            output_path.parent.mkdir(parents=True, exist_ok=True)

            # Check available disk space (basic check)
            if hasattr(os, 'statvfs'):  # Unix systems
                stat = os.statvfs(output_path.parent)
                available_space = stat.f_frsize * stat.f_available
                if available_space < 10 * 1024 * 1024:  # Less than 10MB
                    raise IOError("Insufficient disk space")

            # Save presentation
            prs.save(str(output_path))

        except PermissionError:
            raise IOError(f"Permission denied: Cannot write to {output_path}")
        except OSError as e:
            raise IOError(f"File system error: {e}")


    def _save_speaker_notes_textfile(self, ppt_path: Path, speaker_notes_list: List[Tuple[int, str]]) -> None:
        """Save speaker notes to a text file with spacing between slides"""
        try:
            textfile_path = ppt_path.with_name(ppt_path.stem + "_speaker_notes.txt")
            with open(textfile_path, "w", encoding="utf-8") as f:
                # ✅ Remove duplicates and empty notes
                cleaned_notes = []
                seen = set()

                for slide_num, note in speaker_notes_list:
                    note = note.strip()
                    key = (slide_num, note)

                    if note and key not in seen:
                        seen.add(key)
                        cleaned_notes.append((slide_num, note))

                for slide_number, notes in cleaned_notes:
                    f.write(f"Slide {slide_number}:\n{notes.strip()}\n\n")
            self.logger.info(f"📝 Speaker notes text file saved: {textfile_path}")
        except Exception as e:
            self.logger.error(f"Failed to save speaker notes text file: {e}")

# class CanvasConverter:
#     """Main converter class orchestrating the entire process"""

#     def __init__(self, config: Config = None, log_level: str = 'INFO'):
#         self.config = config or Config()
#         self.logger = setup_logging(log_level)
#         self.extractor = ContentExtractor(self.config, self.logger)
#         self.table_processor = TableProcessor(self.logger)
#         self.ppt_generator = PowerPointGenerator(self.config, self.logger)

#     def convert(self, url: str, output_dir: str = None, filename: str = None) -> Optional[Path]:
#         """Main conversion method with comprehensive error handling"""
#         try:
#             self.logger.info("🚀 Starting ChatGPT Canvas to PowerPoint conversion")

#             # Step 1: Extract content
#             self.logger.info("🔍 Extracting content from canvas...")
#             soup, content = self.extractor.scrape_with_retry(url)

#             if not content:
#                 self.logger.error("❌ Failed to extract content from canvas")
#                 return None

#             # Step 2: Process tables
#             self.logger.info("🧹 Processing tables...")
#             processed_content = self.table_processor.clean_and_validate_tables(soup, content)

#             # Step 3: Generate filename
#             output_path = self._generate_output_path(soup, output_dir, filename)

#             # Step 4: Create PowerPoint
#             self.logger.info("🖼️  Generating PowerPoint presentation...")
#             title = None

#             success = self.ppt_generator.create_enhanced_presentation(
#                 processed_content, output_path, title
#             )

#             if success:
#                 self.logger.info(f"✅ Conversion completed successfully: {output_path}")
#                 return output_path
#             else:
#                 self.logger.error("❌ PowerPoint generation failed")
#                 return None

#         except Exception as e:
#             self.logger.error(f"❌ Conversion failed: {e}")
#             return None

#     def _generate_output_path(self, soup: BeautifulSoup, output_dir: str = None, filename: str = None) -> Path:
#         """Generate safe output path"""
#         if filename:
#             base_name = SafeFilename.sanitize(filename, self.config.max_filename_length)
#         else:
#             # Extract title from page
#             title = self._extract_title(soup)
#             base_name = SafeFilename.sanitize(title, self.config.max_filename_length)

#         if not base_name.endswith('.pptx'):
#             base_name += '.pptx'

#         if output_dir:
#             output_path = Path(output_dir) / base_name
#         else:
#             output_path = Path.cwd() / base_name

#         return SafeFilename.ensure_unique(output_path)

#     def _extract_title(self, soup: BeautifulSoup) -> str:
#         """Extract title from page with multiple strategies"""
#         # Try different title extraction methods
#         title_candidates = [
#             soup.find("h1"),
#             soup.find("title"),
#             soup.find(attrs={"data-testid": "canvas-title"}),
#             soup.find(class_=re.compile("title", re.I)),
#         ]

#         for candidate in title_candidates:
#             if candidate and candidate.get_text(strip=True):
#                 return candidate.get_text(strip=True)

#         return "ChatGPT_Canvas_Export"

#     def batch_convert(self, urls: List[str], output_dir: str = None) -> Dict[str, Optional[Path]]:
#         """Convert multiple canvas URLs"""
#         results = {}

#         for i, url in enumerate(urls, 1):
#             self.logger.info(f"Processing {i}/{len(urls)}: {url}")
#             try:
#                 result = self.convert(url, output_dir)
#                 results[url] = result
#             except Exception as e:
#                 self.logger.error(f"Failed to process {url}: {e}")
#                 results[url] = None

#         return results


# class CLIInterface:
#     """Command-line interface for the converter"""

#     def __init__(self):
#         self.config = Config()
#         self.converter = CanvasConverter(self.config)

#     def run_interactive(self):
#         """Run interactive CLI"""
#         print("🎯 Enhanced ChatGPT Canvas to PowerPoint Converter")
#         print("=" * 60)

#         while True:
#             try:
#                 # Get URL
#                 url = input("\n🔗 Enter ChatGPT canvas share URL (or 'quit' to exit): ").strip()

#                 if url.lower() in ['quit', 'exit', 'q']:
#                     print("👋 Goodbye!")
#                     break

#                 if not url:
#                     print("❌ Please enter a valid URL")
#                     continue

#                 # Get optional parameters
#                 output_dir = input("📁 Output directory (press Enter for current): ").strip() or None
#                 filename = input("📄 Custom filename (press Enter for auto): ").strip() or None

#                 # Convert
#                 result = self.converter.convert(url, output_dir, filename)

#                 if result:
#                     print(f"\n✅ Success! PowerPoint saved to: {result}")

#                     # Ask if user wants to open the file
#                     if self._ask_yes_no("🚀 Open the PowerPoint file now?"):
#                         self._open_file(result)
#                 else:
#                     print("\n❌ Conversion failed. Check the logs above for details.")

#             except KeyboardInterrupt:
#                 print("\n\n👋 Interrupted by user. Goodbye!")
#                 break
#             except Exception as e:
#                 print(f"\n❌ Unexpected error: {e}")

#     def run_batch(self, urls_file: str, output_dir: str = None):
#         """Run batch conversion from file"""
#         try:
#             with open(urls_file, 'r') as f:
#                 urls = [line.strip() for line in f if line.strip() and not line.startswith('#')]

#             if not urls:
#                 print("❌ No valid URLs found in file")
#                 return

#             print(f"🔄 Processing {len(urls)} URLs...")
#             results = self.converter.batch_convert(urls, output_dir)

#             # Summary
#             successful = sum(1 for result in results.values() if result)
#             print(f"\n📊 Batch conversion complete:")
#             print(f"   ✅ Successful: {successful}")
#             print(f"   ❌ Failed: {len(results) - successful}")
#             print(f"   📁 Output directory: {output_dir or 'current directory'}")

#             ppt_paths = [p for p in results.values() if p is not None]
#             if len(ppt_paths) >= 2:
#                 print("🔗 Merging presentations...")
#                 final_merged = merge_presentations(ppt_paths[0], ppt_paths[1:])
#                 print(f"📦 Combined PPT saved at: {final_merged}")

#         except FileNotFoundError:
#             print(f"❌ File not found: {urls_file}")
#         except Exception as e:
#             print(f"❌ Batch processing failed: {e}")

#     def _ask_yes_no(self, question: str) -> bool:
#         """Ask yes/no question"""
#         while True:
#             answer = input(f"{question} (y/n): ").strip().lower()
#             if answer in ['y', 'yes']:
#                 return True
#             elif answer in ['n', 'no']:
#                 return False
#             else:
#                 print("Please enter 'y' or 'n'")

#     def _open_file(self, filepath: Path):
#         """Open file with system default application"""
#         try:
#             if sys.platform.startswith('darwin'):  # macOS
#                 os.system(f'open "{filepath}"')
#             elif sys.platform.startswith('win'):  # Windows
#                 os.startfile(str(filepath))
#             else:  # Linux and others
#                 os.system(f'xdg-open "{filepath}"')
#         except Exception as e:
#             print(f"⚠️  Could not open file automatically: {e}")
#             print(f"📁 File location: {filepath}")


# def main():
#     """Main entry point"""
#     import argparse

#     parser = argparse.ArgumentParser(
#         description="Enhanced ChatGPT Canvas to PowerPoint Converter",
#         formatter_class=argparse.RawDescriptionHelpFormatter,
#         epilog="""
# Examples:
#   python canvas_converter.py                          # Interactive mode
#   python canvas_converter.py -u "https://..."        # Single conversion
#   python canvas_converter.py -b urls.txt -o output/  # Batch conversion
#   python canvas_converter.py -u "https://..." -v     # Verbose logging
#         """
#     )

#     parser.add_argument("-u", "--url", help="ChatGPT canvas URL to convert")
#     parser.add_argument("-o", "--output", help="Output directory")
#     parser.add_argument("-f", "--filename", help="Custom output filename")
#     parser.add_argument("-b", "--batch", help="File containing URLs for batch processing")
#     parser.add_argument("-v", "--verbose", action="store_true", help="Verbose logging")
#     parser.add_argument("--log-level", choices=['DEBUG', 'INFO', 'WARNING', 'ERROR'], 
#                        default='INFO', help="Logging level")
#     parser.add_argument("--merge", action="store_true", help="Merge all generated PPTs into one")


#     args = parser.parse_args()

#     # Set log level
#     log_level = 'DEBUG' if args.verbose else args.log_level

#     # Create CLI interface
#     cli = CLIInterface()
#     cli.converter = CanvasConverter(log_level=log_level)

#     try:
#         if args.batch:
#             # Batch mode
#             cli.run_batch(args.batch, args.output)
#         elif args.url:
#             # Single URL mode
#             result = cli.converter.convert(args.url, args.output, args.filename)
#             if result:
#                 print(f"✅ Success! PowerPoint saved to: {result}")
#             else:
#                 print("❌ Conversion failed")
#                 sys.exit(1)
#         else:
#             # Interactive mode
#             cli.run_interactive()

#     except KeyboardInterrupt:
#         print("\n👋 Interrupted by user")
#     except Exception as e:
#         print(f"❌ Fatal error: {e}")
#         sys.exit(1)

# def merge_presentations(primary_path: Path, additional_paths: List[Path], output_path: Path = None) -> Path:
#     """
#     Merges additional PowerPoint files into a primary one.

#     :param primary_path: Path to the base presentation
#     :param additional_paths: List of presentations to merge in order
#     :param output_path: Path to save the merged presentation
#     :return: Path to the final merged presentation
#     """
#     merged = Presentation(primary_path)

#     for ppt_path in additional_paths:
#         to_merge = Presentation(ppt_path)

#         for slide in to_merge.slides:
#             # Copy slide
#             slide_layout = merged.slide_layouts[1]  # Use content layout
#             new_slide = merged.slides.add_slide(slide_layout)

#             for shape in slide.shapes:
#                 try:
#                     el = shape.element
#                     new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
#                 except Exception as e:
#                     print(f"⚠️ Error copying shape: {e}")

#             # Copy notes
#             try:
#                 notes_text = slide.notes_slide.notes_text_frame.text
#                 if notes_text:
#                     notes = new_slide.notes_slide.notes_text_frame
#                     notes.text = notes_text
#             except Exception:
#                 pass

#     final_path = output_path or primary_path.with_stem(primary_path.stem + "_merged")
#     final_path = SafeFilename.ensure_unique(final_path.with_suffix('.pptx'))
#     merged.save(final_path)
#     return final_path

# if __name__ == "__main__":
#     main()